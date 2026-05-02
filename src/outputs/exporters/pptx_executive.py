"""Executive PowerPoint slide deck. Phase 2 full implementation."""
import tempfile
from pathlib import Path


async def generate(analysis_id: str, version) -> str:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    blank = prs.slide_layouts[6]  # blank layout

    def add_slide(title_text: str, body_text: str):
        slide = prs.slides.add_slide(blank)
        tx_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        tx_title.text_frame.text = title_text
        tx_title.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
        tx_title.text_frame.paragraphs[0].runs[0].font.bold = True

        tx_body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        tx_body.text_frame.word_wrap = True
        tx_body.text_frame.text = body_text
        tx_body.text_frame.paragraphs[0].runs[0].font.size = Pt(18)

    add_slide("Incident Summary", f"Analysis ID: {analysis_id[:8]}")
    add_slide("Root Cause", version.root_cause)
    add_slide("Immediate Resolution", version.workaround)

    actions_text = "\n".join(
        f"• [{a['priority'].upper()}] {a['title']}" for a in version.recommended_actions
    )
    add_slide("Recommended Actions", actions_text)

    out = Path(tempfile.mkdtemp()) / f"{analysis_id}_executive.pptx"
    prs.save(str(out))
    return str(out)
