"""Customer-facing PowerPoint: what happened, what we did, prevention plan."""
import tempfile
from pathlib import Path


async def generate(analysis_id: str, version) -> str:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    blank = prs.slide_layouts[6]

    def add_slide(title_text: str, body_text: str):
        slide = prs.slides.add_slide(blank)
        tx_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        tx_title.text_frame.text = title_text
        tx_title.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
        tx_title.text_frame.paragraphs[0].runs[0].font.bold = True

        tx_body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        tx_body.text_frame.word_wrap = True
        tx_body.text_frame.text = body_text
        tx_body.text_frame.paragraphs[0].runs[0].font.size = Pt(18)

    add_slide("Service Incident Report", "Prepared for customer review")
    add_slide("What Happened", version.root_cause)
    add_slide("What We Did", version.workaround)

    prevention = "\n".join(
        f"• {a['title']}" for a in version.recommended_actions
        if a.get("type") in ("process", "monitoring", "product_improvement")
    ) or "Prevention actions to be determined."
    add_slide("How We're Preventing Recurrence", prevention)

    out = Path(tempfile.mkdtemp()) / f"{analysis_id}_customer.pptx"
    prs.save(str(out))
    return str(out)
